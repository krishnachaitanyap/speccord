#!/usr/bin/env python3
"""Build a downloadable speccord deck (.pptx) mirroring the SVG slides.
Run: python3 build_pptx.py  ->  ../speccord-deck.pptx   (requires python-pptx)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG = RGBColor(0x07, 0x17, 0x1A)
PANEL = RGBColor(0x0F, 0x24, 0x2A)
TEXT = RGBColor(0xD6, 0xEE, 0xF0)
MUTED = RGBColor(0x8F, 0xB2, 0xB6)
SKY = RGBColor(0x22, 0xD3, 0xEE)
EMERALD = RGBColor(0x2D, 0xD4, 0xBF)
VIOLET = RGBColor(0x5E, 0xEA, 0xD4)
AMBER = RGBColor(0x67, 0xE8, 0xF9)
ACCENTS = [SKY, EMERALD, VIOLET, AMBER]
EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, runs, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font="Helvetica Neue"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = runs if isinstance(runs, list) else [runs]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def kicker(s, k, title, accent):
    box(s, 0, 0, 0.13, SH, accent)
    text(s, 0.7, 0.5, 11, 0.5, k, size=15, color=accent, bold=True)
    text(s, 0.7, 0.95, 12, 1.0, title, size=34, color=TEXT, bold=True)


def footer(s, n, total):
    text(s, 0.7, SH - 0.55, 4, 0.4, "speccord", size=13, color=MUTED, bold=True)
    text(s, SW - 1.7, SH - 0.55, 1, 0.4, f"{n} / {total}", size=13, color=MUTED, align=PP_ALIGN.RIGHT)


MONO = "Consolas"
slides = []


def s_title(n, tot):
    s = slide()
    box(s, 0, 2.55, SW, 0.04, SKY)
    box(s, 0, 2.62, SW, 0.04, EMERALD)
    box(s, 0, 2.69, SW, 0.04, VIOLET)
    text(s, 0, 0.8, SW, 1.7, "speccord", size=96, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    text(s, 0, 2.75, SW, 0.6, "the spec is the contract your code is checked against",
         size=22, color=SKY, align=PP_ALIGN.CENTER)
    text(s, 0, 3.9, SW, 0.5, "A spec-driven development CLI — a configurable superset of",
         size=18, color=MUTED, align=PP_ALIGN.CENTER)
    labels = [("spec-kit · generative", SKY), ("speccord · enforcement", EMERALD), ("BMAD · agentic agile", VIOLET)]
    bx = (SW - 3 * 3.2 - 2 * 0.3) / 2
    for i, (lab, c) in enumerate(labels):
        x = bx + i * (3.2 + 0.3)
        b = box(s, x, 4.6, 3.2, 0.6, BG)
        b.line.color.rgb = c
        b.line.width = Pt(2)
        text(s, x, 4.68, 3.2, 0.5, lab, size=15, color=c, bold=True, align=PP_ALIGN.CENTER)
    footer(s, n, tot)


def s_problem(n, tot):
    s = slide()
    kicker(s, "THE PROBLEM", "Specs drift from code. Then they lie.", SKY)
    cards = [("Docs rot", "Hand-written specs go stale the moment code changes. Nobody trusts them."),
             ("Pure-LLM specs hallucinate", "Generate-everything tools invent endpoints, tables, and scopes."),
             ("No enforcement", "Nothing fails the build when the contract and the code disagree.")]
    for i, (h, b) in enumerate(cards):
        x = 0.7 + i * 4.25
        box(s, x, 2.2, 3.9, 3.1, PANEL)
        box(s, x, 2.2, 3.9, 0.09, ACCENTS[i])
        text(s, x + 0.3, 2.5, 3.4, 0.8, h, size=20, color=ACCENTS[i], bold=True)
        text(s, x + 0.3, 3.3, 3.4, 1.8, b, size=15, color=MUTED)
    text(s, 0.7, 5.6, 12, 0.6, "speccord's answer: facts are deterministic — only prose is model-written.",
         size=18, color=TEXT, bold=True)
    footer(s, n, tot)


def s_hybrid(n, tot):
    s = slide()
    kicker(s, "THE IDEA", "Hybrid by construction", EMERALD)
    box(s, 0.7, 2.2, 5.7, 3.2, PANEL)
    text(s, 1.0, 2.4, 5.2, 0.5, "Deterministic (code)", size=19, color=EMERALD, bold=True)
    text(s, 1.0, 3.0, 5.2, 2.3, ["• Parse OpenAPI · migrations · Kafka · security",
                                  "• Diff the contract surface (drift)",
                                  "• Lifecycle state machine + entry gates",
                                  "• Capability resolution from scale"], size=15, color=TEXT)
    box(s, 6.9, 2.2, 5.7, 3.2, PANEL)
    text(s, 7.2, 2.4, 5.2, 0.5, "Model-written (prose only)", size=19, color=VIOLET, bold=True)
    text(s, 7.2, 3.0, 5.2, 2.3, ["• Draft narrative around established facts",
                                  "• Plans, PRDs, stories, reviews",
                                  "• Never invents endpoints/tables/scopes",
                                  "• Grounded in spec + constitution"], size=15, color=TEXT)
    text(s, 0, 5.7, SW, 0.5, "Pass/fail is always code. The model only writes the words.",
         size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    footer(s, n, tot)


def s_lineages(n, tot):
    s = slide()
    kicker(s, "THE SUPERSET", "Three lineages, one CLI", VIOLET)
    rows = [("spec-kit", SKY, "Generative chain", "constitution -> spec -> plan -> tasks -> implement"),
            ("speccord", EMERALD, "Extraction + enforcement", "discover · lifecycle gates · CI drift gate · conform"),
            ("BMAD-METHOD", VIOLET, "Agentic agile", "role personas · scale-adaptive · PRD -> epics -> stories")]
    y = 2.2
    for name, c, role, detail in rows:
        box(s, 0.7, y, 11.9, 1.15, PANEL)
        box(s, 0.7, y, 0.09, 1.15, c)
        text(s, 1.0, y + 0.18, 3.6, 0.6, name, size=21, color=c, bold=True)
        text(s, 1.0, y + 0.66, 3.6, 0.4, role, size=14, color=MUTED)
        text(s, 4.7, y + 0.35, 7.6, 0.5, detail, size=15, color=TEXT, font=MONO)
        y += 1.35
    footer(s, n, tot)


def s_scale(n, tot):
    s = slide()
    kicker(s, "CONFIGURABLE", "One knob: scale (0–4)", AMBER)
    text(s, 0.7, 1.95, 12, 0.4, "Scale sets active phases, enabled roles, and default capabilities — all overridable.",
         size=15, color=MUTED)
    levels = [("0", "prototype", "implement only"), ("1", "small", "+ plan, stories, gate"),
              ("2", "medium", "+ solutioning, PRD, QA"), ("3", "large", "+ analysis, UX"),
              ("4", "enterprise", "+ PO, compliance")]
    w = 2.27
    for i, (lv, nm, add) in enumerate(levels):
        x = 0.7 + i * (w + 0.1)
        c = ACCENTS[i % 4]
        box(s, x, 2.5, w, 2.5, PANEL)
        text(s, x, 2.7, w, 1.0, lv, size=54, color=c, bold=True, align=PP_ALIGN.CENTER)
        text(s, x, 3.75, w, 0.4, nm, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.15, 4.2, w - 0.3, 0.7, add, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, 0.7, 5.5, 12, 0.4, "speccord capabilities", size=17, color=EMERALD, bold=True, font=MONO)
    text(s, 0.7, 5.95, 12, 0.4, "shows what's on, which command each unlocks, and how to change it.",
         size=15, color=MUTED)
    footer(s, n, tot)


def s_phases(n, tot):
    s = slide()
    kicker(s, "THE WORKFLOW", "Four phases → commands", SKY)
    cols = [("ANALYSIS", SKY, ["brief  (analyst)", "discover", "research"]),
            ("PLANNING", EMERALD, ["prd  (pm)", "base draft | new", "feature new · clarify"]),
            ("SOLUTIONING", VIOLET, ["plan  (architect)", "agent architect", "epics & stories"]),
            ("IMPLEMENTATION", AMBER, ["story new (sm)", "implement (dev)", "review (qa)"])]
    w = 2.85
    for i, (ph, c, cmds) in enumerate(cols):
        x = 0.7 + i * (w + 0.2)
        box(s, x, 2.2, w, 3.2, PANEL)
        text(s, x, 2.35, w, 0.5, ph, size=15, color=c, bold=True, align=PP_ALIGN.CENTER)
        text(s, x + 0.25, 3.0, w - 0.4, 2.2, cmds, size=15, color=TEXT, font=MONO)
    text(s, 0.7, 5.9, 12, 0.4, "Always-on: constitution · analyze · checklist · lint · status · advance · gate · agent",
         size=14, color=MUTED)
    footer(s, n, tot)


def s_enforce(n, tot):
    s = slide()
    kicker(s, "THE SPINE", "Enforcement is code, not vibes", EMERALD)
    cards = [("Lifecycle", SKY, "Entry gates: lint, base ref, ACs<->tests, plan+tasks.",
              "Draft -> In Review -> Approved -> ..."),
             ("CI drift gate", VIOLET, "Contract file changed but no spec updated -> build fails.",
              "speccord gate --base origin/main"),
             ("Runtime conform", AMBER, "Re-discovers the live surface, diffs the baseline -> fails on drift.",
              "speccord conform")]
    y = 2.2
    for h, c, body, code in cards:
        box(s, 0.7, y, 11.9, 1.25, PANEL)
        box(s, 0.7, y, 0.09, 1.25, c)
        text(s, 1.0, y + 0.2, 5.5, 0.5, h, size=19, color=c, bold=True)
        text(s, 1.0, y + 0.72, 6.0, 0.5, body, size=14, color=MUTED)
        text(s, 7.2, y + 0.42, 5.2, 0.5, code, size=15, color=TEXT, font=MONO)
        y += 1.45
    footer(s, n, tot)


def s_agents(n, tot):
    s = slide()
    kicker(s, "INTEGRATION", "Plugs into your AI agent", EMERALD)
    text(s, 0.7, 1.95, 12, 0.4, "speccord agent-rules  wires up MCP + a spec-as-contract ruleset for each tool.",
         size=15, color=MUTED, font=MONO)
    cards = [("Claude Code", SKY, ".mcp.json", "CLAUDE.md"),
             ("Cursor", VIOLET, ".cursor/mcp.json", ".cursor/rules/*.mdc"),
             ("Copilot / VS Code", AMBER, ".vscode/mcp.json", "copilot-instructions.md")]
    for i, (nm, c, a, b) in enumerate(cards):
        x = 0.7 + i * 4.25
        box(s, x, 2.45, 3.9, 1.5, PANEL)
        box(s, x, 2.45, 3.9, 0.09, c)
        text(s, x + 0.3, 2.65, 3.4, 0.5, nm, size=18, color=c, bold=True)
        text(s, x + 0.3, 3.15, 3.4, 0.4, a, size=14, color=TEXT, font=MONO)
        text(s, x + 0.3, 3.5, 3.4, 0.4, b, size=14, color=MUTED, font=MONO)
    box(s, 0.7, 4.2, 11.9, 1.4, RGBColor(0x0B, 0x20, 0x26))
    text(s, 1.0, 4.4, 11, 0.4, "MCP tools (speccord mcp):", size=16, color=EMERALD, bold=True)
    text(s, 1.0, 4.8, 11.5, 0.4, "read -> get_base_spec · get_constitution · story_next", size=15, color=TEXT, font=MONO)
    text(s, 1.0, 5.15, 11.5, 0.4, "verify -> analyze · lint · gate · conform    transition -> advance",
         size=15, color=TEXT, font=MONO)
    text(s, 0.7, 5.8, 12, 0.4, "The agent reads the spec, writes code, runs the gates, fixes, then advances. Loop.",
         size=16, color=TEXT, bold=True)
    footer(s, n, tot)


def s_start(n, tot):
    s = slide()
    kicker(s, "TWO MINUTES", "Get started", VIOLET)
    box(s, 0.7, 2.2, 11.9, 3.3, RGBColor(0x0B, 0x20, 0x26))
    lines = [("# existing service", MUTED), ("speccord init --service orders", EMERALD),
             ("speccord discover  &&  speccord base draft", TEXT), ("", TEXT),
             ("# new service", MUTED), ("speccord init --service giftcards --greenfield", EMERALD),
             ('speccord base new --intent "..."', TEXT), ("", TEXT),
             ("# product / BMAD-style", MUTED),
             ("speccord init --pack product  &&  speccord capabilities", EMERALD)]
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.45), Inches(11.3), Inches(2.9))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (ln, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.size = Pt(15)
        r.font.name = MONO
        r.font.color.rgb = c
    text(s, 0.7, 5.7, 12, 0.4, "Read USAGE.md for full tutorials · every pass/fail decision is deterministic.",
         size=15, color=MUTED)
    footer(s, n, tot)


builders = [s_title, s_problem, s_hybrid, s_lineages, s_scale, s_phases, s_enforce, s_agents, s_start]


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    tot = len(builders)
    for i, b in enumerate(builders, 1):
        b(i, tot)
    out = os.path.join(here, "..", "speccord-deck.pptx")
    prs.save(out)
    print("wrote", os.path.abspath(out))


if __name__ == "__main__":
    main()
